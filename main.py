#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
main2.py

目的:
  セキュリティ規則一覧ルールをチェックする CUI ベースの AgentRAG チャットボット

要件（主なポイント）:
  - 起動時に `rule/` の json ファイルを再帰的に読み込む
  - 起動時に `specification/` の pdf/docx/md/txt を読み込み、ChromaDB に格納する
  - Agent A: ドキュメント要約エージェント
  - Agent B: ドキュメント確認エージェント
    - Agent C/D は不要になったため本実装では Agent A/B に集中します
  - LCEL（Runnable 等）を利用してチェーンを組み立てる
  - OpenAI の gpt-4o を ChatOpenAI で呼び出す（OPENAI_API_KEY を利用）

注記:
  - 可読性優先・日本語コメント多め
  - メモリ節約: ドキュメントはチャンクし、検索時は上位 k のみを使用
  - 依存パッケージは README にも書いていますが、`requirements.txt` を参照してください
"""

import os
import sys
import json
import glob
import re
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
import textwrap
from datetime import datetime

# --- 依存パッケージのインポート（利用環境でインストールされている前提） ---
try:
    # LangChain とモデルラッパー
    from langchain.chat_models import ChatOpenAI
    from langchain.schema import HumanMessage, SystemMessage
    from langchain.embeddings import HuggingFaceEmbeddings
    from langchain.vectorstores import Chroma
    from langchain.text_splitter import CharacterTextSplitter
    from langchain.docstore.document import Document
except Exception as e:
    print("必要なライブラリが見つかりません: langchain 等。\n`pip install -r requirements.txt` を実行してください。\nエラー: ", e)
    sys.exit(1)

# LCEL 系のインポート（利用可能なら利用する）
USE_LCEL = True
# 動的 import を使って静的解析エラーの発生を抑え、実行時に利用可能な実装を探す。
import importlib

Runnable = None
RunnablePassthrough = None
try:
    # まず新しい独立パッケージ名を試す
    mod = importlib.import_module("langchain_experimental")
    Runnable = getattr(mod, "Runnable", None)
    RunnablePassthrough = getattr(mod, "RunnablePassthrough", None)
except Exception:
    try:
        # 次に langchain 内の experimental モジュールを試す
        mod2 = importlib.import_module("langchain.experimental.runnable")
        Runnable = getattr(mod2, "Runnable", None)
        RunnablePassthrough = getattr(mod2, "RunnablePassthrough", None)
    except Exception:
        # 利用不可
        Runnable = None
        RunnablePassthrough = None

if Runnable is None or RunnablePassthrough is None:
    USE_LCEL = False

import logging

logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# 環境変数から API キーを取得
OPENAI_API_KEY = os.environ.get("OPENAI_API_KEY")
if not OPENAI_API_KEY:
    print("環境変数 OPENAI_API_KEY が設定されていません。設定してから再実行してください。")
    sys.exit(1)

# --- 設定値 ---
BASE_DIR = Path(__file__).parent
RULE_DIR = BASE_DIR / "rule"
SPEC_DIR = BASE_DIR / "specification"
CHROMA_DIR = BASE_DIR / "chroma_db"
CHROMA_COLLECTION = "specs"

# Embedding モデル名（ローカルで実行できる軽量モデルを使用）
# Chroma の "ローカルの埋め込み" という要件は、ここではローカル HuggingFace 埋め込みを利用して満たす
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"

# RAG 検索時に取り出すドキュメント数
TOP_K = 3


def load_rules_from_dir(rule_dir: Path) -> List[Dict[str, Any]]:
    """rule ディレクトリ配下の全 JSON を再帰的に読み込み、ルールをフラットなリストで返す。

    ルール JSON は配列 または {"rules": [...] } の形式に対応。
    再帰的な構造（子ルールを 'children' などで持つ）もフラット化して返す。
    
    新機能: touitsukijun_r7.json の階層構造 (sections -> subsections -> items) に対応
    """
    rules: List[Dict[str, Any]] = []
    
    def create_rule_entry(item: Dict[str, Any], parent_path: str = "", file_source: str = "") -> Dict[str, Any]:
        """個別ルールエントリを作成"""
        rid = item.get("id") or item.get("rule_id") or item.get("name") or None
        title = item.get("title") or item.get("name") or rid or "unnamed"
        path_label = f"{parent_path}/{title}" if parent_path else title
        
        # content 構築: description + information (あれば)
        content_parts = []
        if item.get("description"):
            content_parts.append(item["description"])
        if item.get("information"):
            content_parts.append(f"\n[詳細情報]\n{item['information']}")
        if item.get("content"):
            content_parts.append(item["content"])
            
        content = "\n".join(content_parts) if content_parts else json.dumps(item, ensure_ascii=False)
        
        return {
            "id": rid,
            "title": title,
            "path": path_label,
            "content": content,
            "type": item.get("type", "未分類"),  # 遵守事項、基本対策事項など
            "source_file": file_source,
            "raw": item,
        }
    
    for path in rule_dir.rglob("*.json"):
        try:
            with open(path, "r", encoding="utf-8") as f:
                data = json.load(f)
        except Exception as e:
            logger.warning(f"ルールファイルを読み込めませんでした: {path} - {e}")
            continue

        file_source = path.name
        
        # 入れ子対応: data が配列か辞書か
        candidates = []
        if isinstance(data, list):
            candidates = data
        elif isinstance(data, dict):
            if "rules" in data and isinstance(data["rules"], list):
                candidates = data["rules"]
            else:
                # 辞書そのものを一つのルール集合とみなす
                candidates = [data]

        # 階層構造の再帰的処理
        def walk(item: Dict[str, Any], parent_path: str = ""):
            # 現在の項目がルールとして追加すべきものかチェック
            has_description = bool(item.get("description") or item.get("content"))
            
            if has_description:
                # ルールエントリとして追加
                entry = create_rule_entry(item, parent_path, file_source)
                rules.append(entry)
            
            # 階層を下って子要素を処理
            current_path = f"{parent_path}/{item.get('title', item.get('id', ''))}" if parent_path else (item.get('title') or item.get('id') or "")
            
            # 新しい階層キー: sections, subsections, items に対応
            for child_key in ("children", "rules", "subrules", "items", "sections", "subsections"):
                if child_key in item and isinstance(item[child_key], list):
                    for child in item[child_key]:
                        if isinstance(child, dict):
                            walk(child, current_path)

        for it in candidates:
            if isinstance(it, dict):
                walk(it)

    logger.info(f"読み込んだルール数: {len(rules)}")
    return rules


def text_from_pdf(path: Path) -> str:
    """シンプルな PDF テキスト抽出。pypdf を利用。ページごとに連結する。"""
    try:
        import pypdf
    except Exception:
        raise RuntimeError("pypdf が必要です。pip install pypdf を実行してください。")
    text_parts = []
    try:
        reader = pypdf.PdfReader(str(path))
        for p in reader.pages:
            txt = p.extract_text() or ""
            text_parts.append(txt)
    except Exception as e:
        logger.warning(f"PDF 読み込み失敗 {path}: {e}")
    return "\n".join(text_parts)


def text_from_docx(path: Path) -> str:
    try:
        import docx
    except Exception:
        raise RuntimeError("python-docx が必要です。pip install python-docx を実行してください。")
    try:
        doc = docx.Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        logger.warning(f"DOCX 読み込み失敗 {path}: {e}")
        return ""


def text_from_xlsx(path: Path) -> str:
    """Excel ファイル(.xlsx)からテキストを抽出"""
    try:
        import openpyxl
    except Exception:
        raise RuntimeError("openpyxl が必要です。pip install openpyxl を実行してください。")
    
    text_parts = []
    try:
        workbook = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        for sheet_name in workbook.sheetnames:
            sheet = workbook[sheet_name]
            text_parts.append(f"\n--- シート: {sheet_name} ---\n")
            
            for row in sheet.iter_rows(values_only=True):
                row_text = []
                for cell in row:
                    if cell is not None:
                        row_text.append(str(cell))
                if row_text:  # 空行でない場合のみ追加
                    text_parts.append("\t".join(row_text))
        workbook.close()
    except Exception as e:
        logger.warning(f"XLSX 読み込み失敗 {path}: {e}")
        return ""
    
    return "\n".join(text_parts)


def text_from_pptx(path: Path) -> str:
    """PowerPoint ファイル(.pptx)からテキストを抽出"""
    try:
        from pptx import Presentation
    except Exception:
        raise RuntimeError("python-pptx が必要です。pip install python-pptx を実行してください。")
    
    text_parts = []
    try:
        prs = Presentation(str(path))
        for i, slide in enumerate(prs.slides, 1):
            text_parts.append(f"\n--- スライド {i} ---\n")
            
            # スライド内の全テキストを抽出
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_parts.append(shape.text.strip())
                
                # 表がある場合のテキスト抽出
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            text_parts.append("\t".join(row_text))
    except Exception as e:
        logger.warning(f"PPTX 読み込み失敗 {path}: {e}")
        return ""
    
    return "\n".join(text_parts)


def load_spec_documents(spec_dir: Path) -> List[Document]:
    """`specification/` 配下のドキュメントを読み込み、langchain Document のリストを返す。

    対応: pdf, docx, xlsx, pptx, md, txt
    メモリ節約: ファイル毎にチャンク分割を行い、最低限のメタデータを付与
    """
    docs: List[Document] = []
    text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
    for path in spec_dir.rglob("*"):
        if path.is_dir():
            continue
        lower = path.suffix.lower()
        try:
            if lower == ".pdf":
                text = text_from_pdf(path)
            elif lower == ".docx":
                text = text_from_docx(path)
            elif lower == ".xlsx":
                text = text_from_xlsx(path)
            elif lower == ".pptx":
                text = text_from_pptx(path)
            elif lower in (".md", ".txt"):
                text = path.read_text(encoding="utf-8", errors="ignore")
            else:
                # 未対応ファイルは無視
                logger.debug(f"未対応ファイル形式をスキップ: {path}")
                continue
        except Exception as e:
            logger.warning(f"ファイル読み込み失敗 {path}: {e}")
            continue

        if not text.strip():
            logger.debug(f"空のファイルをスキップ: {path}")
            continue

        chunks = text_splitter.split_text(text)
        for i, c in enumerate(chunks):
            meta = {"source": str(path), "chunk": i, "file_type": lower}
            docs.append(Document(page_content=c, metadata=meta))

    logger.info(f"読み込んだドキュメントチャンク数: {len(docs)}")
    return docs


def init_chroma(docs: List[Document]) -> Chroma:
    """ChromaDB を初期化して、ドキュメントを格納する。既存コレクションがあれば再利用。

    埋め込みはローカル HuggingFace モデルを使う（軽量モデル推奨）
    """
    # HuggingFace 埋め込みラッパー
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

    # Chroma を初期化（persist_directory を指定して永続化）
    vectordb = Chroma(persist_directory=str(CHROMA_DIR), collection_name=CHROMA_COLLECTION, embedding_function=embeddings)

    # 既存が空の場合は追加
    try:
        # Chroma が提供する API により既存数を確認
        existing = vectordb._collection.count() if hasattr(vectordb, "_collection") else None
    except Exception:
        existing = None

    if existing in (None, 0):
        if docs:
            logger.info("Chroma にドキュメントを追加します...")
            vectordb.add_documents(docs)
            vectordb.persist()
    else:
        logger.info("既存の Chroma コレクションを利用します。新規追加は行いません。")

    return vectordb


def init_rules_chroma(rules: List[Dict[str, Any]], force_rebuild: bool = False) -> Optional[Chroma]:
    """ルールを専用の ChromaDB コレクションに格納する（オプション機能）。

    Note: この機能は現在無効化されています。メモリ内でのルール処理のみ使用します。

    Args:
        rules: ルールのリスト
        force_rebuild: True の場合、既存コレクションを削除して再構築
        
    Returns:
        ルール専用の ChromaDB インスタンス（現在は None を返す）
    """
    logger.info("ルールのベクトルDB化は無効化されています。メモリ内処理のみ使用します。")
    return None
    if not rules:
        return None
        
    # HuggingFace 埋め込みラッパー
    embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)

    # ルール専用コレクション
    rule_vectordb = Chroma(persist_directory=str(CHROMA_DIR), collection_name="rules", embedding_function=embeddings)

    # 強制再構築の場合は既存データを削除
    if force_rebuild:
        try:
            rule_vectordb._collection.delete()
            logger.info("既存のルールコレクションを削除しました")
            # 新しいコレクションを作成
            rule_vectordb = Chroma(persist_directory=str(CHROMA_DIR), collection_name="rules", embedding_function=embeddings)
        except Exception as e:
            logger.warning(f"既存コレクション削除でエラー: {e}")

    # 既存が空の場合は追加
    try:
        existing = rule_vectordb._collection.count() if hasattr(rule_vectordb, "_collection") else None
    except Exception:
        existing = None

    if existing in (None, 0) or force_rebuild:
        logger.info("ルールを Chroma に追加します...")
        # ルールを Document として変換
        rule_docs = []
        for rule in rules:
            # ルール内容をページコンテンツに
            content = rule.get('content', '')
            if not content:
                continue
                
            # メタデータに詳細情報を含める
            metadata = {
                "rule_id": rule.get('id', ''),
                "title": rule.get('title', ''),
                "type": rule.get('type', ''),
                "path": rule.get('path', ''),
                "source_file": rule.get('source_file', ''),
                "doc_type": "rule"  # ドキュメントと区別するためのマーカー
            }
            
            # 長いルールの場合はチャンク分割
            if len(content) > 1500:
                text_splitter = CharacterTextSplitter(chunk_size=1000, chunk_overlap=200)
                chunks = text_splitter.split_text(content)
                for i, chunk in enumerate(chunks):
                    chunk_metadata = metadata.copy()
                    chunk_metadata["chunk"] = i
                    rule_docs.append(Document(page_content=chunk, metadata=chunk_metadata))
            else:
                rule_docs.append(Document(page_content=content, metadata=metadata))
        
        if rule_docs:
            # バッチサイズ制限を回避するために分割して追加
            batch_size = 100  # ChromaDBの制限より小さい値に設定
            total_docs = len(rule_docs)
            logger.info(f"ルール {total_docs} チャンクを {batch_size} ずつバッチで ChromaDB に追加します")
            
            for i in range(0, total_docs, batch_size):
                batch = rule_docs[i:i + batch_size]
                logger.info(f"バッチ {i//batch_size + 1}/{(total_docs + batch_size - 1)//batch_size}: {len(batch)} チャンクを追加中...")
                try:
                    rule_vectordb.add_documents(batch)
                except Exception as e:
                    logger.warning(f"バッチ {i//batch_size + 1} の追加に失敗: {e}")
                    # さらに小さなバッチで再試行
                    smaller_batch_size = 50
                    for j in range(i, min(i + batch_size, total_docs), smaller_batch_size):
                        smaller_batch = rule_docs[j:j + smaller_batch_size]
                        try:
                            rule_vectordb.add_documents(smaller_batch)
                            logger.info(f"小バッチ {j//smaller_batch_size + 1}: {len(smaller_batch)} チャンク追加成功")
                        except Exception as e2:
                            logger.error(f"小バッチ {j//smaller_batch_size + 1} も失敗: {e2}")
            
            rule_vectordb.persist()
            logger.info(f"ルール {total_docs} チャンクの ChromaDB 追加が完了しました")
    else:
        logger.info("既存のルール ChromaDB コレクションを利用します。")

    return rule_vectordb


def retrieve_related_docs(vectordb: Chroma, query: str, k: int = TOP_K, include_rules: bool = False) -> List[Document]:
    """関連ドキュメントを検索する。オプションでルールも含めることができる。

    Args:
        vectordb: メインの仕様ドキュメント用 ChromaDB
        query: 検索クエリ
        k: 取得する文書数
        include_rules: True の場合、ルールコレクションからも検索（デフォルト無効）
    
    Returns:
        関連するドキュメントのリスト
    """
    docs = []
    
    # メインドキュメントから検索
    try:
        main_docs = vectordb.similarity_search(query, k=k)
        docs.extend(main_docs)
    except Exception as e:
        logger.warning(f"メインドキュメント検索エラー: {e}")
    
    # ルールからも検索（オプション、デフォルトは無効）
    if include_rules:
        try:
            rule_embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
            rule_vectordb = Chroma(persist_directory=str(CHROMA_DIR), collection_name="rules", embedding_function=rule_embeddings)
            rule_docs = rule_vectordb.similarity_search(query, k=k//2)  # 半分はルールから
            docs.extend(rule_docs)
        except Exception as e:
            logger.warning(f"ルール検索エラー（スキップします）: {e}")
    
    return docs[:k]  # 上位k件に制限


# --- LCEL ベースの簡易 Wrapper（利用可能なら Runnable を用いる） ---
def make_chat_model() -> ChatOpenAI:
    """ChatOpenAI を作成。モデル名は gpt-4o を指定する。"""
    # ChatOpenAI は環境変数 OPENAI_API_KEY を参照する
    llm = ChatOpenAI(model_name="gpt-4o", temperature=0.0)
    return llm


def agent_a_summarize(llm: ChatOpenAI, rule_text: str, docs: List[Document]) -> str:
    """Agent A: ドキュメント要約エージェント

    - 複数ドキュメントとルールを受け取り、要約を返す
    - 完全性を重視する（重要事項の漏れがないように指示）
    """
    # 取得テキストを簡潔にまとめる（上位 k 件のみを渡す）
    context = "\n\n---関連ドキュメント---\n"
    for d in docs[:TOP_K]:
        src = d.metadata.get("source") if d.metadata else "<unknown>"
        context += f"[source: {src}]\n{d.page_content}\n\n"

    system_prompt = (
        "あなたは優秀なドキュメント要約者です。以下のセキュリティルールを読み、"
        "関連ドキュメントの内容を完全性を保って要約してください。重要な条件、要件、及び検証ポイントを箇条書きで示してください。"
        "\n\n重要: 回答は必ず日本語で行ってください。出力の本文は日本語で記載し、読みやすい箇条書きを心がけてください。"
    )
    messages = [SystemMessage(content=system_prompt), HumanMessage(content=f"ルール:\n{rule_text}\n\nドキュメントコンテキスト:\n{context}")]
    resp = llm(messages)
    return resp.content


def agent_b_check(llm: ChatOpenAI, rule_summary: str, rule_raw: Dict[str, Any], docs: List[Document]) -> Dict[str, Any]:
    """Agent B: ドキュメント確認エージェント

    - ルールに対して、システム（docs）が従っているか評価する
    - 出力は構造化 JSON で返す（result, evidence, details）
    """
    context = "\n\n".join([f"[src:{d.metadata.get('source')}]\n{d.page_content}" for d in docs[:TOP_K]])

    # 厳密な JSON 出力を促すプロンプト（スキーマと例を明示）
    strict_prompt = """
あなたは技術的な評価者です。以下のルール要約と元ルールを読み、与えられたシステム情報がそのルールに従っているかを評価してください。

出力は厳密な JSON のみを返してください。余計な説明や追加テキストは一切書かず、必ず純粋な JSON テキストだけを返してください（コードフェンスや説明を含めないでください）。

JSON スキーマ例:
{
    "result": "〇|△|×",            // 判定
    "evidence": [                   // 推奨: 配列形式
        {"source": "ファイル名や識別子", "excerpt": "抜粋テキスト..."}
    ],
    "details": "追加の説明(任意)"
}

重要: JSON のキー名は英語のままにし、値や説明文は日本語で記載してください。
"""

    init_human = HumanMessage(content=f"ルール要約:\n{rule_summary}\n\n元ルール(raw):\n{json.dumps(rule_raw, ensure_ascii=False)}\n\nドキュメントコンテキスト:\n{context}")

    messages = [SystemMessage(content=strict_prompt), init_human]

    # 最初の回答を取得
    resp = llm(messages)
    # モデルに JSON 形式で答えるよう指示したが、念のためパースを試みる
    text = resp.content

    # デバッグ用: 生出力をログに残す（短縮版）
    logger.debug("Agent B raw output (head 1000 chars): %s", text[:1000].replace('\n', '\\n'))

    # 万が一のため、失敗した生出力をファイルへ追記するユーティリティ
    def _save_model_output(rule_id: str, content: str):
        try:
            logs_dir = BASE_DIR / "logs"
            logs_dir.mkdir(exist_ok=True)
            ts = datetime.utcnow().strftime("%Y%m%dT%H%M%SZ")
            fname = logs_dir / f"agent_b_output_{ts}_{str(rule_id)[:60].replace(' ', '_')}.log"
            with open(fname, "w", encoding="utf-8") as lf:
                lf.write("--- RAW MODEL OUTPUT ---\n")
                lf.write(content)
            logger.info("モデル出力をログに保存しました: %s", fname)
        except Exception as e:
            logger.debug("モデル出力ログ保存に失敗しました: %s", e)

    # 保存は任意（環境変数で無効化可能）
    if os.environ.get("SAVE_MODEL_OUTPUT", "1") != "0":
        try:
            _save_model_output(rule_raw.get("id") or rule_raw.get("title") or "unknown", text)
        except Exception:
            pass

    parsed = None
    # まず素直に JSON としてデコードを試みる
    try:
        parsed = json.loads(text)
    except Exception:
        parsed = None

    # もしパース失敗したらモデルに再試行を促す（最大2回）
    retries = 0
    while parsed is None and retries < 2:
        retries += 1
        logger.info("JSON パース失敗: モデルへ再試行を行います (試行 %d)。", retries)
        followup = (
            "前の回答は有用でしたが、要求された通り厳密な JSON のみで出力されていませんでした。"
            "以下の JSON スキーマに厳密に合わせ、純粋な JSON テキストのみを出力してください。"
            "\n\nスキーマ: {\"result\":\"〇|△|×\", \"evidence\": [ {\"source\":..., \"excerpt\":...} ], \"details\": \"任意の文字列\" }"
            "\n\n元の出力を参照して、上記スキーマにマッピングして JSON のみを返してください。"
        )
        follow_messages = [SystemMessage(content=strict_prompt), init_human, HumanMessage(content=followup + "\n\n前の出力:\n" + text)]
        try:
            resp2 = llm(follow_messages)
            text2 = resp2.content
            logger.debug("Agent B retry raw output (head 1000): %s", text2[:1000].replace('\n', '\\n'))
            # まず素直に JSON としてデコードを試みる
            try:
                parsed = json.loads(text2)
                text = text2
                break
            except Exception:
                # 次に波括弧ブロックを抽出
                m2 = re.search(r"(\{[\s\S]*\})", text2)
                if m2:
                    try:
                        parsed = json.loads(m2.group(1))
                        text = m2.group(1)
                        break
                    except Exception:
                        parsed = None
                # 最後に簡易変換を試す
                t2 = text2.replace("'", '"')
                t2 = re.sub(r",\s*([}\]])", r"\1", t2)
                t2 = re.sub(r'([\{,\s])(\w+)\s*:', r'\1"\2":', t2)
                try:
                    parsed = json.loads(t2)
                    text = t2
                    break
                except Exception:
                    parsed = None
        except Exception as e:
            logger.debug("モデル再試行中に例外: %s", e)
            parsed = None

    # ここまでで parsed が None ならオリジナル text を用いてフォールバック処理へ
    if parsed is None:
        # JSON の部分文字列を抜き出す試み
        # 1) 最初の波括弧ブロックを抽出して試す
        m = re.search(r"(\{[\s\S]*\})", text)
        if m:
            candidate = m.group(1)
            try:
                parsed = json.loads(candidate)
            except Exception:
                parsed = None

        # 2) シングルクォートをダブルクォートに変換、末尾の余分なカンマを削除、未引用キーに引用付与を試す
        if parsed is None:
            t2 = text.replace("'", '"')
            t2 = re.sub(r",\s*([}\]])", r"\1", t2)
            t2 = re.sub(r'([\{,\s])(\w+)\s*:', r'\1"\2":', t2)
            try:
                parsed = json.loads(t2)
            except Exception:
                parsed = None

    # 最終的にパースできなければヒューリスティック抽出へ
    if parsed is None:
        logger.warning("モデルの出力を JSON としてパースできませんでした。ヒューリスティック抽出を試みます。")

        def _heuristic_parse(text: str) -> Dict[str, Any]:
            out: Dict[str, Any] = {}
            # result (期待値: 〇/△/× または O/X など)
            m = re.search(r"['\"]?result['\"]?\s*[:：]\s*['\"]?([^\"',}\n\r]+)", text, re.IGNORECASE)
            if m:
                out["result"] = m.group(1).strip().strip('"\'')
            else:
                m2 = re.search(r"\b(〇|△|×|O|X|o|x)\b", text)
                if m2:
                    out["result"] = m2.group(1)

            # evidence: try to extract the value after evidence key (allow multiline)
            m_e = re.search(r"['\"]?evidence['\"]?\s*[:：]\s*([\"'])(.*?)\1", text, re.IGNORECASE | re.DOTALL)
            if m_e:
                out["evidence"] = m_e.group(2).strip()
            else:
                m_e2 = re.search(r"evidence\s*[:：\-]\s*(.+)$", text, re.IGNORECASE | re.DOTALL)
                if m_e2:
                    out["evidence"] = m_e2.group(1).strip()

            # details
            m_d = re.search(r"['\"]?details['\"]?\s*[:：]\s*([\"'])(.*?)\1", text, re.IGNORECASE | re.DOTALL)
            if m_d:
                out["details"] = m_d.group(2).strip()
            else:
                m_d2 = re.search(r"details\s*[:：\-]\s*(.+)$", text, re.IGNORECASE | re.DOTALL)
                if m_d2:
                    out["details"] = m_d2.group(1).strip()

            # もし何も抽出できなければ全体を evidence として格納
            if not out.get("evidence") and text:
                out["evidence"] = text.strip()

            # 最低限の result を設定
            if "result" not in out:
                out["result"] = "△"

            return out

        parsed = _heuristic_parse(text)

    # 正規化: evidence は list[ {source, excerpt} ] の形にする
    def _build_evidence_list(evidence_field, docs_list: List[Document]):
        evs = []
        # もしモデルが文字列を返してきたら、そのまま root evidence として docs の抜粋を付与
        if not evidence_field:
            # フォールバック: docs の先頭から抜粋を作る
            for d in docs_list[:TOP_K]:
                evs.append({"source": d.metadata.get("source"), "excerpt": d.page_content[:400].strip()})
            return evs

        if isinstance(evidence_field, str):
            # モデルのフリーテキストをそのまま一つの根拠とする
            evs.append({"source": "(model-output)", "excerpt": evidence_field})
            # さらに docs から抜粋を付与
            for d in docs_list[:TOP_K]:
                evs.append({"source": d.metadata.get("source"), "excerpt": d.page_content[:300].strip()})
            return evs

        # リスト形式が期待される場合
        if isinstance(evidence_field, list):
            for item in evidence_field:
                if isinstance(item, dict):
                    src = item.get("source") or item.get("file") or item.get("path") or "(unknown)"
                    exc = item.get("excerpt") or item.get("text") or json.dumps(item, ensure_ascii=False)
                    evs.append({"source": src, "excerpt": exc[:400].strip()})
                else:
                    evs.append({"source": "(model-output)", "excerpt": str(item)[:400]})
            return evs

        # それ以外の型は文字列化して格納
        evs.append({"source": "(model-output)", "excerpt": str(evidence_field)[:400]})
        return evs

    parsed_evidence = _build_evidence_list(parsed.get("evidence"), docs)
    parsed["evidence_normalized"] = parsed_evidence
    return parsed


# Agent C はユーザの要望により本実装では除外しました。


# Agent D (ダブルチェック) は要件から除外されたため定義は削除しました。


def retrieve_related_docs_simple(vectordb: Chroma, query: str, k: int = TOP_K) -> List[Document]:
    """簡易 RAG 用: query を埋め込み検索し、上位 k のドキュメントを返す（後方互換性用）"""
    retriever = vectordb.as_retriever(search_kwargs={"k": k})
    results = retriever.get_relevant_documents(query)
    return results


def format_b_result(b_result: Dict[str, Any]) -> str:
    """Agent B の構造化結果を日本語の整形テキストに変換する"""
    lines: List[str] = []
    res = b_result.get("result") or b_result.get("status") or "△"
    lines.append(f"判定: {res}")
    # 詳細説明
    details = b_result.get("details") or b_result.get("detail") or b_result.get("notes")
    if details:
        lines.append("\n説明:")
        if isinstance(details, str):
            lines.append(details)
        else:
            lines.append(json.dumps(details, ensure_ascii=False, indent=2))

    # 根拠（正規化済み）
    evs = b_result.get("evidence_normalized") or []
    if evs:
        lines.append("\n根拠 (参照文書と抜粋):")
        for i, e in enumerate(evs, 1):
            src = e.get("source") or "(unknown)"
            excerpt = e.get("excerpt") or ""
            # 抜粋は改行削除して短めに
            excerpt_clean = excerpt.replace("\n", " ").strip()
            if len(excerpt_clean) > 1200:
                excerpt_clean = excerpt_clean[:1200].rstrip() + " ..."
            lines.append(f"  {i}. source: {src}")
            # indent excerpt
            wrapped = textwrap.fill(excerpt_clean, width=100, subsequent_indent='     ')
            lines.append(textwrap.indent(wrapped, '     '))

    # モデル本体の自由テキスト evidence がある場合
    if b_result.get("evidence") and not evs:
        lines.append("\nモデル出力（根拠）:")
        lines.append(str(b_result.get("evidence")))

    return "\n".join(lines)


 


def find_rule_by_query(rules: List[Dict[str, Any]], query: str) -> Optional[Dict[str, Any]]:
    """ルール一覧から query を元にルールを検索する。ID/パス/タイトル/本文の部分一致で最初のマッチを返す。"""
    q = query.strip().lower()
    # まず ID にマッチ
    for r in rules:
        if r.get("id") and str(r.get("id")).lower() == q:
            return r
    # その他のフィールドで部分マッチ
    for r in rules:
        if q in (r.get("title") or "").lower() or q in (r.get("path") or "").lower() or q in (r.get("content") or "").lower():
            return r
    return None


def interactive_loop(rules: List[Dict[str, Any]], vectordb: Chroma):
    """CUI ベースの簡易チャットループ"""
    llm = make_chat_model()

    help_text = (
        "コマンド一覧:\n"
        "  help                      ヘルプ表示\n"
        "  list                      読み込んだルール一覧の一部を表示\n"
        "  show <query>              ルールを表示（id/title の部分一致）\n"
        "  check <query>             指定したルールに対してシステムが従っているか評価（A->B の順）\n"
        "  showfull <summary|b>      直近のチェックで保存された項目の全文表示\n"
        "  ask <自由テキスト>        システム情報に関する RAG 質問\n"
        "  quit                      終了\n"
    )

    print("AgentRAG チャットボット (CUI)。help と入力してください。\n")
    # 直近の run で作られた出力を保存（全文表示用）
    last_store: Dict[str, Any] = {"summary": None, "b": None}

    def print_section(title: str, content: str, max_len: int = 1200):
        """見やすいセクション表示。

        - セクション見出しを強調表示
        - テキストは段落単位で折り返しを行い、読みやすくする
        - JSON 文字列らしい場合は折り返しせずそのまま出力する
        """
        sep = "=" * 80
        print("\n" + sep)
        print(f"{title}")
        print(sep)
        if content is None:
            print("(なし)\n")
            return

        # プレビュー長を超える場合は省略の目印を付与
        display_text = content
        if isinstance(content, str) and len(content) > max_len:
            display_text = content[:max_len].rstrip() + "\n...（全文は 'showfull' コマンドで表示可）"

        # JSONらしい出力はそのまま表示
        if isinstance(display_text, str) and display_text.strip().startswith(("{", "[")):
            print(display_text)
        else:
            # 段落ごとに折り返して表示（空行で段落分割）
            if isinstance(display_text, str):
                paras = [p.strip() for p in display_text.split("\n\n") if p.strip()]
                for p in paras:
                    wrapped = textwrap.fill(p, width=100)
                    print(wrapped)
                    print()
            else:
                print(str(display_text))
        print(sep + "\n")
    while True:
        try:
            cmd = input("> ").strip()
        except (EOFError, KeyboardInterrupt):
            print("\n終了します。")
            break

        if not cmd:
            continue
        if cmd == "help":
            print(help_text)
            continue
        if cmd == "list":
            for i, r in enumerate(rules[:50], 1):
                print(f"{i}. id={r.get('id')} title={r.get('title')} path={r.get('path')}")
            continue
        if cmd.startswith("show "):
            q = cmd[len("show "):].strip()
            r = find_rule_by_query(rules, q)
            if not r:
                print("ルールが見つかりませんでした。部分文字列で検索してみてください。")
            else:
                print("--- ルール ---")
                print(f"id: {r.get('id')}")
                print(f"title: {r.get('title')}")
                print(f"path: {r.get('path')}")
                print("content:")
                print(r.get("content"))
            continue

        if cmd.startswith("check "):
            q = cmd[len("check "):].strip()
            r = find_rule_by_query(rules, q)
            if not r:
                print("ルールが見つかりません。別のクエリを試してください。\n(例: ルールの一部の語句や id を入力)\n")
                continue

            print(f"選択されたルール: {r.get('title')} (path: {r.get('path')})")

            # Agent A: 要約（関連ドキュメントを検索して渡す）
            rule_text = r.get("content") or ""
            related = retrieve_related_docs(vectordb, rule_text, k=TOP_K)
            print("[Agent A] ルールと関連ドキュメントから要約を作成しています...")
            summary = agent_a_summarize(llm, rule_text, related)
            last_store["summary"] = summary
            print_section("Agent A - 要約プレビュー", summary)

            # Agent B: 確認
            print("[Agent B] ドキュメントがルールに従っているか評価しています...")
            b_result = agent_b_check(llm, summary, r.get("raw", {}), related)
            last_store["b"] = b_result
            # B の判定表示: プレビュー + 根拠の参照文書表記
            # 整形テキストで表示する
            b_preview_text = format_b_result(b_result)
            print_section("Agent B - 判定（プレビュー）", b_preview_text)
            # 根拠一覧を見やすく表示
            evs = b_result.get("evidence_normalized") or []
            if evs:
                print("根拠 (参照文書と抜粋):")
                for i, e in enumerate(evs, 1):
                    src = e.get("source") or "(unknown)"
                    excerpt = e.get("excerpt") or ""
                    print(f"  {i}. source: {src}")
                    # 抜粋は改行を整形して折り返し表示
                    excerpt_clean = excerpt.replace("\n", " ").strip()
                    if len(excerpt_clean) > 1000:
                        excerpt_clean = excerpt_clean[:1000].rstrip() + " ..."
                    wrapped = textwrap.fill(excerpt_clean, width=100, subsequent_indent='     ')
                    print(textwrap.indent(wrapped, '     '))
                    print()
            else:
                print("(根拠情報はありません)")
            # 補足: 簡易のアクション提案
            result_symbol = b_result.get("result", "△")
            if result_symbol == "〇" or result_symbol == "O" or result_symbol == "o":
                print("補足: 判定は '従っている' と見なされます。必要に応じて関連資料を参照してください。\n")
            elif result_symbol == "×" or result_symbol == "X" or result_symbol == "x":
                print("補足: 判定は '従っていない' です。優先的な対応（修正／設定変更等）が必要です。詳細は関連資料を参照してください。\n")
            else:
                print("補足: 判定は '△'（追加確認が必要）です。関連箇所のログや設定ファイルを追加で提供してください。\n")

            # 評価フロー完了（Agent C は除外）
            print("評価フローが完了しました。必要に応じて 'showfull summary' や 'showfull b' で全文を表示できます。")
            continue

        if cmd.startswith("showfull "):
            what = cmd[len("showfull "):].strip()
            if what not in ("summary", "b"):
                print("'showfull' の引数は summary|b のいずれかを指定してください。")
                continue
            val = last_store.get(what)
            if val is None:
                print(f"まだ '{what}' の出力がありません。先に 'check <query>' を実行してください。")
                continue
            title_map = {"summary": "Agent A - 要約（全文）", "b": "Agent B - 判定（全文）"}
            # dict の場合は整形テキストに変換して出力
            if isinstance(val, str):
                content = val
            else:
                if what == "b":
                    content = format_b_result(val)
                else:
                    content = json.dumps(val, ensure_ascii=False, indent=2)
            print_section(title_map.get(what, what), content, max_len=10_000)
            continue

        if cmd.startswith("ask "):
            q = cmd[len("ask "):].strip()
            # RAG 質問: ドキュメントから上位 TOP_K を引いて LLM に渡す
            docs = retrieve_related_docs(vectordb, q, k=TOP_K)
            context = "\n\n".join([f"[src:{d.metadata.get('source')}]\n{d.page_content}" for d in docs])
            system = "あなたはシステム情報の検索アシスタントです。ユーザの質問に、関連するドキュメントを参照して簡潔に答えてください。"
            messages = [SystemMessage(content=system), HumanMessage(content=f"質問: {q}\n\n参照文書:\n{context}")]
            resp = llm(messages)
            print(resp.content)
            continue

        if cmd in ("quit", "exit", "q"):
            print("終了します。")
            break

        print("不明なコマンドです。help を表示してください。")


def main():
    # ルール読み込み
    rules = load_rules_from_dir(RULE_DIR)

    # ドキュメント読み込み
    docs = load_spec_documents(SPEC_DIR)

    # Chroma 初期化・インデックス作成
    vectordb = init_chroma(docs)
    
    # ルール用ChromaDB初期化（ルールのクロスリファレンス機能）
    rule_vectordb = init_rules_chroma(rules)

    # インタラクティブループ開始
    interactive_loop(rules, vectordb)


if __name__ == "__main__":
    main()
